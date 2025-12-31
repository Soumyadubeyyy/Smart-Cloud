# main.py
import os
import io
import json
import uuid
import fitz
import docx
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv
from fastapi import FastAPI, File, UploadFile, HTTPException, Depends, Query
from fastapi.middleware.cors import CORSMiddleware
import google.generativeai as genai
from PIL import Image
from pydantic import BaseModel, Field
from sqlalchemy.orm import Session  
from sqlalchemy import text
from sqlalchemy import func
import secrets
from datetime import datetime, timedelta, timezone
from fastapi.responses import RedirectResponse
from passlib.context import CryptContext

import database
import models
from supabase_client import supabase
from auth import get_current_user
from hashing import generate_hash
from embedding import generate_embedding
from supabase_storage import upload_file_to_storage, get_download_url, delete_file_from_storage


load_dotenv()
pwd_context = CryptContext(schemes=["argon2"], deprecated="auto")
app = FastAPI(title="Smart Cloud Backend API")

# Create database tables & get db session
models.Base.metadata.create_all(bind=database.engine)
def get_db():
    db = database.SessionLocal()
    try:
        yield db
    finally:
        db.close()


app_status = {
    "gemini_status": "INITIALIZING",
    "startup_error": None
}

# GEMINI API CONFIGURATION
try:
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise ValueError("GEMINI_API_KEY not found in .env file.")
    genai.configure(api_key=api_key)
    app_status["gemini_status"] = "OK"
except Exception as e:
    error_message = f"Error configuring Gemini library: {e}"
    app_status["gemini_status"] = "ERROR"
    app_status["startup_error"] = str(e)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# extraction functions
def extract_text_from_pdf(content_bytes: bytes) -> str:
    text = ""
    with fitz.open(stream=content_bytes, filetype="pdf") as doc:
        for page in doc: text += page.get_text()
    return text

def extract_text_from_docx(content_bytes: bytes) -> str:
    text = ""
    doc = docx.Document(io.BytesIO(content_bytes))
    for para in doc.paragraphs: text += para.text + "\n"
    return text

def extract_text_from_pptx(content_bytes: bytes) -> str:
    text = ""
    prs = Presentation(io.BytesIO(content_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"): text += shape.text + "\n"
    return text

def extract_text_from_xlsx(content_bytes: bytes) -> str:
    text = ""
    xls = pd.ExcelFile(io.BytesIO(content_bytes))
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name, header=None)
        text += df.to_string(index=False, header=False) + "\n"
    return text


def get_gemini_analysis(text: str) -> dict:
    """Gets category, summary, and embedding using the genai.GenerativeModel & Voyage Multimodal embedding."""
    if app_status["gemini_status"] != "OK":
        raise HTTPException(
            status_code=503,
            detail={"message": "Service unavailable: Gemini library failed to initialize.", "startup_error": app_status["startup_error"]}
        )
    if not text:
        return {"category": "Other", "summary": "No text content found."}
    
    try:
        # Task 1: Generate Category and Summary
        
        model = genai.GenerativeModel("gemini-2.5-flash")
        prompt = f"""
        Analyze the following text. Respond ONLY with a valid JSON object with keys "category" and "summary".
        1. "category": A single category from [
    "Work", "Finance", "Education", "Legal", "Health", "Entertainment", 
    "Travel", "Household", "Career", "Creative", "Technical", "Shopping", 
    "Social", "Personal", "Archival", "Other"
].
        2. "summary": A concise, one-paragraph summary.

        Text: '{text[:15000]}'
        """
        response = model.generate_content(prompt)
        cleaned_response = response.text.strip().replace("```json", "").replace("```", "")
        result_json = json.loads(cleaned_response)
        category = result_json.get("category", "Other")
        summary = result_json.get("summary", "Summary not available.")

        return {"category": category, "summary": summary}

    except Exception as e:
        print(f"Error during Gemini API call: {e}")
        raise HTTPException(status_code=500, detail=f"Gemini API Error: {e}")
    
#Pydantic models for Auth
class UserCredentials(BaseModel):
    email: str
    password: str

class FileResponse(BaseModel):
    id: int
    original_filename: str
    category: str | None
    summary: str | None
    upload_date: datetime
    file_size_bytes: int

    class Config:
        from_attributes = True

# Auth endpoints : SignUp & SignIn & SignOut
@app.post("/auth/signup")
async def signup(credentials: UserCredentials):
    try:
        response = supabase.auth.sign_up({"email": credentials.email, "password": credentials.password})
        return {"message": "Signup successful, please check your email to verify.", "data": response}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/auth/login")
async def login(credentials: UserCredentials):
    try:
        response = supabase.auth.sign_in_with_password({"email": credentials.email, "password": credentials.password})
        return response
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))
    
@app.post("/auth/signout")
async def signout(user = Depends(get_current_user)):
    """
    Signs out the current user by invalidating their session token.
    """
    try:
        # We need the user's token, which our auth dependency gives us.
        # However, the supabase-py library needs the raw token string from the header.
        # A robust way is to re-implement the token extraction here.
        # For simplicity in this step, we assume the frontend can pass it.
        # A more advanced auth dependency would provide the raw token.
        
        # The supabase-py sign_out function invalidates the token on the server.
        supabase.auth.sign_out() 
        
        return {"message": "User signed out successfully."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

NEAR_DUPLICATE_THRESHOLD = 97.5
# File upload & analysis endpoint
@app.post("/upload-and-analyze/")
async def upload_and_analyze_file(file: UploadFile = File(...),force: bool = Query(False, description="Set to true to bypass the near-duplicate check."),user = Depends(get_current_user), db: Session = Depends(get_db)):
    
    # Duplicacy check
    file_hash = await generate_hash(file)
    existing_file = db.query(models.File).filter(models.File.user_id == user.id, models.File.file_hash == file_hash).first()
    if existing_file:
        raise HTTPException(status_code=409, detail="File with this content already exists.")
    
    
    content_bytes = await file.read()
    await file.seek(0) # Reset file pointer after reading
    
    category, summary, new_embedding, extracted_text = None, None, None, ""
    content_type = file.content_type

    if content_type.startswith('image/'):
        print("Image file detected. Generating embedding only...")
        image = Image.open(io.BytesIO(content_bytes))
        new_embedding = generate_embedding(content=image)
        category = "Image"
        summary = file.filename
    elif content_type.startswith(('video/', 'audio/')):
        category = content_type.split('/')[0].capitalize()
        summary = file.filename
        new_embedding = None
    else:
        if content_type == 'application/pdf':
            extracted_text = extract_text_from_pdf(content_bytes)
        elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            extracted_text = extract_text_from_docx(content_bytes)
        elif content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            extracted_text = extract_text_from_pptx(content_bytes)
        elif content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            extracted_text = extract_text_from_xlsx(content_bytes)
        elif content_type.startswith('text/'):
            extracted_text = content_bytes.decode('utf-8', errors='ignore')
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type for analysis: {content_type}")
        
        if extracted_text:
            new_embedding = generate_embedding(content=extracted_text)
            analysis = get_gemini_analysis(extracted_text)
            category = analysis["category"]
            summary = analysis["summary"]
        else:
            category = "Other"
            summary = "No text content could be extracted."
            new_embedding = None

    if new_embedding and not force:
        # Near-duplicate check using cosine similarity
        query = text("""
            SELECT original_filename, (1 - (embedding <=> :query_embedding)) * 100 AS similarity
            FROM files WHERE user_id = :user_id ORDER BY embedding <=> :query_embedding LIMIT 1
        """)
        result = db.execute(query, {"query_embedding": str(new_embedding), "user_id": str(user.id)}).first()
        if result and result.similarity > NEAR_DUPLICATE_THRESHOLD:
            return {
                "status": "near_duplicate_found",
                "message": "This file is very similar to an existing file.",
                "similar_to_filename": result.original_filename,
                "similarity_score": f"{result.similarity:.2f}%"
            }
        
    file_extension = os.path.splitext(file.filename)[1]
    stored_filename = f"{user.id}/{uuid.uuid4()}{file_extension}"
    
    if not upload_file_to_storage(file, stored_filename):
        raise HTTPException(status_code=500, detail="Failed to upload file to cloud storage.")
        
    new_file_record = models.File(
        user_id=user.id,
        original_filename=file.filename,
        stored_filename=stored_filename,
        file_hash=file_hash,
        file_size_bytes=file.size,
        mime_type=file.content_type,
        category=category,
        summary=summary,
        embedding=new_embedding
    )
    db.add(new_file_record)
    db.commit()
    db.refresh(new_file_record)

    return {"status": "success", "message": "File uploaded successfully!", "file_id": new_file_record.id}
    

@app.get("/health")
def health_check():
    return app_status

@app.post("/generate-hash/")
async def generate_hash_endpoint(file: UploadFile = File(...)):
    """
    Receives a file, generates its SHA-256 hash, and returns it.
    """
    file_hash = await generate_hash(file)
    return {
        "filename": file.filename,
        "sha256_hash": file_hash
    }

@app.get("/files", response_model=list[FileResponse])
async def list_files(user = Depends(get_current_user), db: Session = Depends(get_db)):
    files = db.query(models.File).filter(models.File.user_id == user.id).order_by(models.File.upload_date.desc()).all()
    return files

# Download link endpoint
@app.get("/files/{file_id}/download")
async def get_file_download_link(file_id: int, user = Depends(get_current_user), db: Session = Depends(get_db)):
    """
    Generates a secure, temporary download link for a specific file.
    """
    file_record = db.query(models.File).filter(models.File.id == file_id, models.File.user_id == user.id).first()
    if not file_record:
        raise HTTPException(status_code=404, detail="File not found or you do not have permission to access it.")
    stored_filename = file_record.stored_filename
    download_url = get_download_url(stored_filename)
    if not download_url:
        raise HTTPException(status_code=500, detail="Could not generate download link.")
    return {"download_url": download_url}

# Delete file endpoint
@app.delete("/files/{file_id}")
async def delete_file(file_id: int, user = Depends(get_current_user), db: Session = Depends(get_db)):
    """
    Deletes a specific file and its database record after verifying ownership.
    """
    file_record = db.query(models.File).filter(models.File.id == file_id, models.File.user_id == user.id).first()
    if not file_record:
        raise HTTPException(status_code=404, detail="File not found or you do not have permission to delete it.")
    stored_filename = file_record.stored_filename
    was_deleted_from_storage = delete_file_from_storage(stored_filename)
    if not was_deleted_from_storage:
        raise HTTPException(status_code=500, detail="Failed to delete file from storage.")

    try:
        db.delete(file_record)
        db.commit()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"File deleted from storage, but failed to delete database record: {e}")

    return {"message": f"File '{stored_filename}' was deleted successfully."}


# File sharing logic and endpoints
class ShareLinkCreate(BaseModel):
    password: str | None = Field(default=None, min_length=8, max_length=72)
    expires_in_days: int | None = None

@app.post("/files/{file_id}/share", status_code=201)
async def create_share_link(
    file_id: int, 
    options: ShareLinkCreate, 
    user = Depends(get_current_user), 
    db: Session = Depends(get_db)
):
    """Creates a secure, shareable link for one of the user's files."""
    file_record = db.query(models.File).filter(models.File.id == file_id, models.File.user_id == user.id).first()
    if not file_record:
        raise HTTPException(status_code=404, detail="File not found.")

    token = secrets.token_urlsafe(32)
    password_hash = pwd_context.hash(options.password) if options.password else None
    expires_at = datetime.now(timezone.utc) + timedelta(days=options.expires_in_days) if options.expires_in_days else None

    new_link = models.ShareLink(
        token=token,
        file_id=file_id,
        created_by_user_id=user.id,
        password_hash=password_hash,
        expires_at=expires_at
    )
    db.add(new_link)
    db.commit()
    
    share_url = f"http://localhost:8000/share/{token}"
    return {"share_url": share_url}


@app.get("/share/{token}")
async def access_shared_file(token: str, db: Session = Depends(get_db)):
    """Public endpoint to access a shared file via a token."""
    link_record = db.query(models.ShareLink).filter(models.ShareLink.token == token).first()

    if not link_record:
        raise HTTPException(status_code=404, detail="Share link not found or has been revoked.")

    if link_record.expires_at and link_record.expires_at < datetime.now(timezone.utc):
        raise HTTPException(status_code=410, detail="This share link has expired.")

    if link_record.password_hash:
        raise HTTPException(status_code=401, detail="This link is password protected.")

    file_record = db.query(models.File).filter(models.File.id == link_record.file_id).first()
    if not file_record:
        raise HTTPException(status_code=404, detail="The linked file no longer exists.")

    download_url = get_download_url(file_record.stored_filename)
    if not download_url:
        raise HTTPException(status_code=500, detail="Could not generate download link.")
    
    return RedirectResponse(url=download_url)

#Smart Search
SEARCH_THRESHOLD = 0.2

@app.get("/search", response_model=list[FileResponse])
async def smart_search(query: str, user = Depends(get_current_user), db: Session = Depends(get_db)):
    """
    Performs a semantic search for the user's files based on a query.
    """
    if not query:
        return []
    print(f"Generating embedding for search query: '{query}'")
    query_embedding = generate_embedding(content=query)
    print(query_embedding[:10])
    if not query_embedding:
        raise HTTPException(status_code=500, detail="Could not generate embedding for the search query.")

    sql_query = text("""
        SELECT *
        FROM files
        WHERE user_id = :user_id AND (1 - (embedding <=> :query_embedding)) > :threshold
        ORDER BY embedding <=> :query_embedding
        LIMIT 5
    """)
    
    results = db.execute(
        sql_query,
        {
            "query_embedding": str(query_embedding),
            "user_id": str(user.id),
            "threshold": SEARCH_THRESHOLD
        }
    ).fetchall()
    
    return results

#Analytics pydantic models & endpoint

class CategoryCount(BaseModel):
    category: str
    count: int

class DashboardResponse(BaseModel):
    total_files: int
    total_storage_mb: float
    recent_files: list[FileResponse]  
    category_counts: list[CategoryCount]
    data_version: str

class DashboardStatus(BaseModel):
    status: str
    


@app.get("/dashboard", response_model=DashboardResponse | DashboardStatus)
def get_dashboard_stats(
    current_version: str | None = Query(None),
    user = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    stats_query = db.query(
        func.count(models.File.id), 
        func.max(models.File.upload_date)
    ).filter(models.File.user_id == user.id).first()
    
    count = stats_query[0] or 0
    last_upload = stats_query[1]
    new_version = f"{count}_{last_upload}"
    
    if current_version == new_version:
        return {"status": "unchanged"}

    total_files = count # We already have the count from the fast query above

    total_size_bytes = db.query(func.sum(models.File.file_size_bytes))\
        .filter(models.File.user_id == user.id).scalar() or 0
    total_size_mb = round(total_size_bytes / (1024 * 1024), 2)

    # 1. Get Raw Database Objects
    recent_files_db = db.query(models.File)\
        .filter(models.File.user_id == user.id)\
        .order_by(models.File.upload_date.desc())\
        .limit(5)\
        .all()
    
    # 2. THIS IS THE MISSING FIX: Convert to Pydantic Models manually
    recent_files = [FileResponse.model_validate(f) for f in recent_files_db]

    category_counts = db.query(models.File.category, func.count(models.File.id))\
        .filter(models.File.user_id == user.id)\
        .group_by(models.File.category)\
        .all()
    
    categories = []
    for cat, count_val in category_counts:
        categories.append({"category": cat or "Uncategorized", "count": count_val})
    
    return {
        "total_files": total_files,
        "total_storage_mb": total_size_mb,
        "recent_files": recent_files,  # Now sending clean Pydantic objects
        "category_counts": categories,
        "data_version": new_version
    }

@app.get("/")
def read_root():
    return {"status": "Welcome to the Smart Storage API!"}